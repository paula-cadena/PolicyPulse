import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import sys
import os
from scipy import stats
import joblib

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor



sys.path.append(os.path.dirname(__file__))
from utils import logger

# Setting plotting style
plt.style.use('seaborn-v0_8')
sns.set_palette("husl")

def create_summary_statistics(model_data):
    """Create summary statistics table"""
    logger.info("Creating summary statistics")
    
    # Selecting key variables for summary stats
    stats_vars = ['precipitation', 'precip_mean', 'precip_std', 'precip_anomaly', 
                 'credit_portfolio', 'savings_deposits']
    
    available_vars = [var for var in stats_vars if var in model_data.columns]
    
    summary_stats = model_data[available_vars].describe().round(2)
    
    # Saving summary statistics to CSV
    summary_stats.to_csv('outputs/tables/summary_statistics.csv')
    
    logger.info("Summary statistics saved")
    return summary_stats

def plot_climate_distribution(model_data):
    """Plot distribution of climate variables"""
    logger.info("Creating climate distribution plots")
    
    fig, axes = plt.subplots(2, 2, figsize=(15, 12))
    
    # Precipitation distribution
    if 'precipitation' in model_data.columns:
        axes[0,0].hist(model_data['precipitation'].dropna(), bins=50, alpha=0.7, edgecolor='black')
        axes[0,0].set_xlabel('Monthly Precipitation (mm)')
        axes[0,0].set_ylabel('Frequency')
        axes[0,0].set_title('Distribution of Monthly Precipitation')
    
    # Precipitation anomalies
    if 'precip_anomaly' in model_data.columns:
        axes[0,1].hist(model_data['precip_anomaly'].dropna(), bins=50, alpha=0.7, edgecolor='black', color='orange')
        axes[0,1].axvline(0, color='red', linestyle='--', alpha=0.8)
        axes[0,1].set_xlabel('Precipitation Anomaly (mm)')
        axes[0,1].set_ylabel('Frequency')
        axes[0,1].set_title('Distribution of Precipitation Anomalies')
    
    # Monthly precipitation patterns
    if 'precipitation' in model_data.columns and 'month' in model_data.columns:
        monthly_avg = model_data.groupby('month')['precipitation'].mean()
        axes[1,0].plot(monthly_avg.index, monthly_avg.values, marker='o', linewidth=2)
        axes[1,0].set_xlabel('Month')
        axes[1,0].set_ylabel('Average Precipitation (mm)')
        axes[1,0].set_title('Seasonal Precipitation Patterns')
        axes[1,0].set_xticks(range(1, 13))
    
    # Extreme precipitation events
    if 'extreme_precip' in model_data.columns:
        extreme_counts = model_data['extreme_precip'].value_counts()
        axes[1,1].pie(extreme_counts.values, labels=['Normal', 'Extreme'], autopct='%1.1f%%', 
                     colors=['lightblue', 'red'])
        axes[1,1].set_title('Extreme Precipitation Events')
    
    plt.tight_layout()
    plt.savefig('outputs/figures/climate_distributions.png', dpi=300, bbox_inches='tight')
    plt.close()
    
    logger.info("Climate distribution plots saved")

def plot_financial_outcomes(model_data):
    """Plot distribution of financial outcomes"""
    logger.info("Creating financial outcomes plots")
    
    fig, axes = plt.subplots(1, 2, figsize=(15, 6))
    
    # Credit portfolio distribution
    if 'credit_portfolio' in model_data.columns:
        axes[0].hist(np.log1p(model_data['credit_portfolio'].dropna()), bins=50, alpha=0.7, edgecolor='black')
        axes[0].set_xlabel('Log(Credit Portfolio + 1)')
        axes[0].set_ylabel('Frequency')
        axes[0].set_title('Distribution of Credit Portfolio (log scale)')
    
    # Savings deposits distribution
    if 'savings_deposits' in model_data.columns:
        axes[1].hist(np.log1p(model_data['savings_deposits'].dropna()), bins=50, alpha=0.7, edgecolor='black', color='green')
        axes[1].set_xlabel('Log(Savings Deposits + 1)')
        axes[1].set_ylabel('Frequency')
        axes[1].set_title('Distribution of Savings Deposits (log scale)')
    
    plt.tight_layout()
    plt.savefig('outputs/figures/financial_distributions.png', dpi=300, bbox_inches='tight')
    plt.close()
    
    logger.info("Financial outcomes plots saved")

def plot_climate_financial_relationship(model_data):
    """Plot relationship between climate variables and financial outcomes"""
    logger.info("Creating climate-financial relationship plots")
    
    fig, axes = plt.subplots(2, 2, figsize=(15, 12))
    
    # Precipitation vs Credit Portfolio
    if 'precipitation' in model_data.columns and 'credit_portfolio' in model_data.columns:
        sample_data = model_data.sample(min(1000, len(model_data)), random_state=42)
        axes[0,0].scatter(sample_data['precipitation'], np.log1p(sample_data['credit_portfolio']), 
                         alpha=0.5, s=10)
        axes[0,0].set_xlabel('Monthly Precipitation (mm)')
        axes[0,0].set_ylabel('Log(Credit Portfolio + 1)')
        axes[0,0].set_title('Precipitation vs Credit Portfolio')
        
        # Add trend line
        z = np.polyfit(sample_data['precipitation'].dropna(), 
                      np.log1p(sample_data['credit_portfolio'].dropna()), 1)
        p = np.poly1d(z)
        axes[0,0].plot(sample_data['precipitation'], p(sample_data['precipitation']), "r--", alpha=0.8)
    
    # Precipitation Anomaly vs Credit Portfolio
    if 'precip_anomaly' in model_data.columns and 'credit_portfolio' in model_data.columns:
        sample_data = model_data.sample(min(1000, len(model_data)), random_state=42)
        axes[0,1].scatter(sample_data['precip_anomaly'], np.log1p(sample_data['credit_portfolio']), 
                         alpha=0.5, s=10, color='orange')
        axes[0,1].axvline(0, color='red', linestyle='--', alpha=0.8)
        axes[0,1].set_xlabel('Precipitation Anomaly (mm)')
        axes[0,1].set_ylabel('Log(Credit Portfolio + 1)')
        axes[0,1].set_title('Precipitation Anomaly vs Credit Portfolio')
    
    # Extreme precipitation impact
    if 'extreme_precip' in model_data.columns and 'credit_portfolio' in model_data.columns:
        extreme_effect = model_data.groupby('extreme_precip')['credit_portfolio'].mean()
        axes[1,0].bar(['Normal', 'Extreme'], extreme_effect.values, color=['lightblue', 'red'], alpha=0.7)
        axes[1,0].set_ylabel('Average Credit Portfolio')
        axes[1,0].set_title('Impact of Extreme Precipitation on Credit')
        
        # Add value labels on bars
        for i, v in enumerate(extreme_effect.values):
            axes[1,0].text(i, v, f'{v:,.0f}', ha='center', va='bottom')
    
    # Feature importance plot
    try:
        feature_importance_credit = pd.read_csv('outputs/tables/feature_importance_credit_portfolio.csv')
        top_features = feature_importance_credit.head(10)
        axes[1,1].barh(top_features['feature'], top_features['importance'], color='purple', alpha=0.7)
        axes[1,1].set_xlabel('Feature Importance')
        axes[1,1].set_title('Top 10 Features for Credit Portfolio Prediction')
        axes[1,1].invert_yaxis()
    except FileNotFoundError:
        axes[1,1].text(0.5, 0.5, 'Feature importance data not available', 
                      ha='center', va='center', transform=axes[1,1].transAxes)
    
    plt.tight_layout()
    plt.savefig('outputs/figures/climate_financial_relationships.png', dpi=300, bbox_inches='tight')
    plt.close()
    
    logger.info("Climate-financial relationship plots saved")

def create_policy_implications_table(model_data, results):
    """Create table with policy implications"""
    logger.info("Creating policy implications table")
    
    policy_insights = []
    
    # Calculating key statistics for policy recommendations
    if 'extreme_precip' in model_data.columns and 'credit_portfolio' in model_data.columns:
        normal_credit = model_data[model_data['extreme_precip'] == 0]['credit_portfolio'].mean()
        extreme_credit = model_data[model_data['extreme_precip'] == 1]['credit_portfolio'].mean()
        credit_impact_pct = ((extreme_credit - normal_credit) / normal_credit * 100)
        
        policy_insights.append({
            'metric': 'Credit Portfolio during Extreme Precipitation',
            'normal_conditions': f"${normal_credit:,.0f}",
            'extreme_conditions': f"${extreme_credit:,.0f}",
            'impact': f"{credit_impact_pct:+.1f}%",
            'policy_implication': 'Financial resilience programs needed in climate-vulnerable regions'
        })
    
    if 'precip_std' in model_data.columns:
        high_variability = model_data['precip_std'].quantile(0.75)
        low_variability = model_data['precip_std'].quantile(0.25)
        
        policy_insights.append({
            'metric': 'Precipitation Variability (Standard Deviation)',
            'normal_conditions': f"{low_variability:.1f} mm",
            'extreme_conditions': f"{high_variability:.1f} mm", 
            'impact': f"{(high_variability/low_variability - 1)*100:+.0f}% higher variability",
            'policy_implication': 'Regions with high climate variability need adaptive financial products'
        })
    
    # Creating policy table
    policy_df = pd.DataFrame(policy_insights)
    policy_df.to_csv('outputs/tables/policy_implications.csv', index=False)
    
    logger.info("Policy implications table saved")
    return policy_df

def create_powerpoint_presentation(model_data, results):
    """Create PowerPoint presentation with analysis results"""
    
    # Creating presentation object
    prs = Presentation()
    
    # Function to set professional styling
    def set_slide_styling(slide):
        """Set professional styling for slide elements"""
        # Title styling
        if slide.shapes.title:
            title = slide.shapes.title
            title.text_frame.paragraphs[0].font.size = Pt(30)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Content styling
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            text_frame = content.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)

    def add_centered_image(slide, image_path, title_text=None):
        """Add an image centered on the slide with proper sizing"""
        if not os.path.exists(image_path):
            logger.warning(f"Image not found: {image_path}")
            return False
            
        try:
            # Slide dimensions
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            
            # Image positioning
            left_margin = Inches(0.5)
            right_margin = Inches(0.5)
            top_margin = Inches(1.5)
            bottom_margin = Inches(1.0)
            
            # Available space for image
            available_width = slide_width - left_margin - right_margin
            available_height = slide_height - top_margin - bottom_margin
            
            # Adding title
            title_box = slide.shapes.add_textbox(left_margin, Inches(0.5), 
                                            available_width, Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add image with maximum size constraints
            pic = slide.shapes.add_picture(image_path, left_margin, top_margin, 
                                        width=available_width)
            
            # In case image is too tall, scale by height instead
            if pic.height > available_height:
                pic.height = available_height
                # Re-center horizontally
                pic.left = (slide_width - pic.width) / 2
            
            logger.info(f"Added image: {os.path.basename(image_path)}")
            return True
            
        except Exception as e:
            logger.error(f"Error adding image {image_path}: {e}")
            return False

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Climate Impact on Financial Outcomes"
    subtitle.text = "Analysis of Precipitation Effects on Municipal Finance in Colombia"
    
    # Style title slide
    title.text_frame.paragraphs[0].font.size = Pt(36)
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Slide 2: Introduction and Problem Statement
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title.text = "Introduction & Problem Statement"
    content.text = (
        "Climate change poses significant risks to economic stability\n\n"
        "Research Question: How does climate variability affect financial outcomes at the municipal level in Colombia?\n\n"
        "Objective: Quantify the relationship between precipitation patterns and financial indicators\n\n"
        "Scope: Municipal-level analysis across Colombia using public data"
    )
    set_slide_styling(slide2)
    
    # Slide 3: Data Sources
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide3.shapes.title
    content = slide3.placeholders[1]
    
    title.text = "Data Sources"
    content.text = (
        "Climate Data - IDEAM:\n"
        "Monthly Precipitation Indices (2000-2024)\n"
        "42 municipalities across Colombia\n"
        "Source: Instituto de Hidrología, Meteorología y Estudios Ambientales\n\n"
        "Financial Data - Superintendencia Financiera:\n"
        "Credit Portfolios and Savings Deposits by Municipality\n"
        "Monthly data across multiple years\n"
        "National coverage of financial institutions\n\n"
        "Geographic Data:\n"
        "Colombian Municipality Reference Database\n"
        "Standardized municipality identifiers"
    )
    set_slide_styling(slide3)
    
    # Slide 4: Methodology
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide4.shapes.title
    content = slide4.placeholders[1]
    
    title.text = "Methodology"
    content.text = (
        "Data Processing Pipeline:\n"
        "1. Data Integration & Cleaning\n"
        "2. Municipality Matching (Fuzzy String Matching to improve data quality)\n"
        "3. Feature Engineering for Climate Variability\n"
        "4. Statistical Analysis & Machine Learning\n\n"
        "Analytical Approach:\n"
        "Panel data regression analysis\n"
        "Random Forest for feature importance\n"
        "Climate anomaly calculation\n"
        "Temporal aggregation and alignment"
    )
    set_slide_styling(slide4)
    
    # Slide 5: Climate Data Overview
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_image(slide5, 'outputs/figures/climate_distributions.png', 
                      "Climate Data: Precipitation Patterns")
    
    # Slide 6: Financial Data Overview
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_image(slide6, 'outputs/figures/financial_distributions.png',
                      "Financial Data Distribution")
    
    # Slide 7: Key Relationships
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_image(slide7, 'outputs/figures/climate_financial_relationships.png',
                      "Climate-Financial Relationships")
    
    # Slide 8: Dataset Overview
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide8.shapes.title
    content = slide8.placeholders[1]
    
    title.text = "Dataset Overview"
    
    # Calculate key statistics
    total_records = len(model_data)
    unique_municipalities = model_data['municipality_code'].nunique()
    date_range = f"{model_data['date'].min().strftime('%Y-%m')} to {model_data['date'].max().strftime('%Y-%m')}"
    avg_precipitation = model_data['precipitation'].mean()
    avg_credit = model_data['credit_portfolio'].mean()
    
    overview_text = (
        f"Analysis Coverage:\n"
        f"Total Records: {total_records:,}\n"
        f"Municipalities: {unique_municipalities}\n"
        f"Time Period: {date_range}\n\n"
        f"Key Averages:\n"
        f"Monthly Precipitation: {avg_precipitation:.1f} mm\n"
        f"Credit Portfolio: ${avg_credit:,.0f} COP\n\n"
        f"Data Quality:\n"
        f"Municipality matching: 42 cities\n"
        f"Temporal alignment: Monthly frequency\n"
        f"Coverage: Colombian municipalities with IDEAM monitoring stations"
    )
    
    content.text = overview_text
    set_slide_styling(slide8)
    
    # Slide 9: Regression Results - Credit Portfolio
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide9.shapes.title
    content = slide9.placeholders[1]
    
    title.text = "Regression Analysis: Credit Portfolio"
    
    credit_text = (
        "OLS Regression Results:\n"
        "• R-squared: 0.015 (model explains 1.5% of variance)\n"
        "• Observations: 48,177\n"
        "• Precipitation coefficient: 5.507e+10 (p < 0.001)\n"
        "• Precipitation anomaly: -5.498e+10 (p < 0.001)\n\n"
        "Key Findings:\n"
        "• Higher average precipitation associated with higher credit portfolios\n"
        "• Precipitation anomalies show inverse relationship\n"
        "• Month and quarter effects not statistically significant\n"
        "• Controlling for year effects (banking system expansion)\n\n"
        "Interpretation:\n"
        "Municipalities with higher typical precipitation levels tend to have\n"
        "larger credit portfolios, suggesting climate conditions influence\n"
        "financial activity and credit availability."
    )
    
    content.text = credit_text
    set_slide_styling(slide9)
    
    # Slide 10: Regression Results - Savings Deposits
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide10.shapes.title
    content = slide10.placeholders[1]
    
    title.text = "Regression Analysis: Savings Deposits"
    
    savings_text = (
        "OLS Regression Results:\n"
        "• R-squared: 0.014 (model explains 1.4% of variance)\n"
        "• Observations: 48,177\n"
        "• Precipitation coefficient: 2.957e+10 (p < 0.001)\n"
        "• Precipitation anomaly: -2.949e+10 (p < 0.001)\n\n"
        "Key Findings:\n"
        "• Similar pattern to credit portfolios but smaller coefficients\n"
        "• Precipitation anomalies negatively affect savings\n"
        "• Statistical significance maintained across specifications\n"
        "• Consistent with credit portfolio findings\n\n"
        "Interpretation:\n"
        "Savings behavior also responds to climate conditions, though\n"
        "the effects are slightly smaller than for credit portfolios."
    )
    
    content.text = savings_text
    set_slide_styling(slide10)
    
    # Slide 11: Machine Learning Insights
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide11.shapes.title
    content = slide11.placeholders[1]
    
    title.text = "ML - Random Forest Feature Importance"
    
    # Read feature importance data
    try:
        credit_importance = pd.read_csv('outputs/tables/feature_importance_credit_portfolio.csv')
        savings_importance = pd.read_csv('outputs/tables/feature_importance_savings_deposits.csv')
        
        ml_text = "Random Forest Results - Top Predictors:\n"
        
        ml_text += "Credit Portfolio:\n"
        for _, row in credit_importance.head(3).iterrows():
            feature_name = row['feature'].replace('_', ' ').title()
            ml_text += f"• {feature_name}: {row['importance']:.3f}\n"
        
        ml_text += "\nSavings Deposits:\n"
        for _, row in savings_importance.head(3).iterrows():
            feature_name = row['feature'].replace('_', ' ').title()
            ml_text += f"• {feature_name}: {row['importance']:.3f}\n"
        
        ml_text += "\nKey Insights:\n"
        ml_text += "Local precipitation norms are strongest predictors (65-66%)\n"
        ml_text += "Time trend (year) captures financial system expansion (26-28%)\n"
        ml_text += "Current precipitation and anomalies have smaller but consistent effects\n"
        ml_text += "Extreme weather indicators show minimal direct impact\n\n"
        ml_text += "Interpretation:\n"
        ml_text += "Long-term climate patterns matter more than short-term fluctuations\n"
        ml_text += "in explaining municipal financial outcomes."
    except Exception as e:
        ml_text = f"Feature importance data not available. Error: {e}"
    
    content.text = ml_text
    set_slide_styling(slide11)
    
    
    # Slide 13: Policy Implications
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide13.shapes.title
    content = slide13.placeholders[1]
    
    title.text = "Policy Implications"
    content.text = (
        "1. Climate-Resilient Financial Products\n"
        "   Weather-indexed insurance for agriculture and small businesses\n"
        "   Flexible repayment schedules tied to climate events\n"
        "2. Regional Targeting Strategy\n"
        "   Focus financial inclusion programs on climate-vulnerable regions\n"
        "   Develop early warning systems for financial institutions\n"
        "   Regional adaptation funds for climate-resilient investments\n\n"
        "3. Financial Inclusion Enhancement\n"
        "   Mobile banking for remote climate-affected areas\n"
        "   Climate-resilient microfinance products\n"
        "   Financial literacy programs focused on climate risks"
    )
    set_slide_styling(slide13)
    
    # Slide 14: Implementation Challenges
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide14.shapes.title
    content = slide14.placeholders[1]
    
    title.text = "Implementation Challenges & Solutions"
    content.text = (
        "Financial Sector Concentration\n"
        "Challenge: Majority of financial services from few private institutions\n"
        "Solution: Public-private partnerships and regulatory incentives\n\n"
        "Geographic Access Limitations\n"
        "Challenge: Climate-vulnerable regions have less financial access\n"
        "Solution: Mobile banking expansion and agent banking networks\n\n"
        "Informality Challenges\n"
        "Challenge: High informality limits formal financial product effectiveness\n"
        "Solution: Gradual formalization programs and hybrid products\n\n"
    )
    set_slide_styling(slide14)
    
    # Slide 15: Analysis Limitations
    slide15 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide15.shapes.title
    content = slide15.placeholders[1]
    
    title.text = "Analysis Limitations & Future Improvements"
    content.text = (
        "Data Quality Constraints:\n"
        "Municipality-level climate data completeness and accuracy\n"
        "Financial services concentration in urban centers\n"
        "Temporal coverage variations across municipalities\n\n"
        "Methodological Considerations:\n"
        "Correlation analysis, not causal inference\n"
        "Potential omitted variable bias (economic activity, infrastructure)\n"
        "Aggregation effects in larger municipalities\n\n"
        "Scope for Enhancement:\n"
        "Additional variables: Credit types, distance to urban centers\n"
        "Expanded climate indicators: Temperature, drought indices\n"
        "Causal methods: Instrumental variables, natural experiments\n\n"
    )
    set_slide_styling(slide15)
    
    # Slide 16: Conclusions and Next Steps
    slide16 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide16.shapes.title
    content = slide16.placeholders[1]
    
    title.text = "Conclusions and Next Steps"
    content.text = (
        "Analytical Contributions:\n"
        "Demonstrated measurable climate-finance relationship in Colombia\n"
        "Developed reproducible analytical pipeline with public data\n"
        "Applied both statistical and machine learning approaches\n"
        "Policy Relevance:\n"
        "Quantified climate impacts on municipal financial outcomes\n"
        "Identified key predictors for targeted interventions\n"
        "Highlighted implementation challenges and solutions\n\n"
        "Recommended Next Steps:\n"
        "1. Expand analysis with additional variables and municipalities\n"
        "2. Develop real-time monitoring dashboard\n"
        "3. Conduct causal analysis with natural experiments\n"
        "4. Engage stakeholders for policy implementation\n"
        "5. Explore sector-specific climate finance impacts"
    )
    set_slide_styling(slide16)

    # Save presentation
    output_path = 'docs/slides.pptx'
    prs.save(output_path)
    logger.info(f"Enhanced presentation saved to: {output_path}")
    logger.info(f"Total slides created: {len(prs.slides)}")
    
    return True

def main():
    """Main evaluation and visualization function"""
    logger.info("Starting evaluation and visualization process")
    
    # Creating output directories
    os.makedirs('outputs/figures', exist_ok=True)
    os.makedirs('outputs/tables', exist_ok=True)
    os.makedirs('docs', exist_ok=True)
    
    # Loading modeling data
    model_data = pd.read_csv('data_processed/modeling_dataset_cleaned.csv')
    model_data['date'] = pd.to_datetime(model_data['date'])
    logger.info(f"Loaded modeling data: {len(model_data)} records")

    
    # Loading results
    results = {}
    try:
        feature_importance_credit = pd.read_csv('outputs/tables/feature_importance_credit_portfolio.csv')
        results['credit_portfolio'] = {'feature_importance': feature_importance_credit}
    except FileNotFoundError:
        logger.warning("Model results not found, proceeding with descriptive analysis only")
    
    # Creating all visualizations and tables
    summary_stats = create_summary_statistics(model_data)
    plot_climate_distribution(model_data)
    plot_financial_outcomes(model_data)
    plot_climate_financial_relationship(model_data)
    policy_table = create_policy_implications_table(model_data, results)
    
    # Creating PowerPoint presentation
    presentation_success = create_powerpoint_presentation(model_data, results)
    
    # Adding key findings to log
    logger.info("\n" + "="*50)
    logger.info("KEY FINDINGS SUMMARY")
    logger.info("="*50)
    
    if 'precipitation' in model_data.columns:
        logger.info(f"• Average monthly precipitation: {model_data['precipitation'].mean():.1f} mm")
        logger.info(f"• Precipitation variability (std): {model_data['precipitation'].std():.1f} mm")
    
    if 'credit_portfolio' in model_data.columns:
        logger.info(f"• Average credit portfolio: ${model_data['credit_portfolio'].mean():,.0f}")
    
    if 'extreme_precip' in model_data.columns:
        extreme_events = model_data['extreme_precip'].sum()
        logger.info(f"• Extreme precipitation events: {extreme_events} ({extreme_events/len(model_data)*100:.1f}% of observations)")
    
    if presentation_success:
        logger.info("• PowerPoint presentation created successfully in docs/slides.pptx")
    
    logger.info("Evaluation and visualization completed successfully")

if __name__ == "__main__":
    main()